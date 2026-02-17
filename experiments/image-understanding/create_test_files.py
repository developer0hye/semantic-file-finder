"""
테스트용 DOCX/PPTX 생성 스크립트

이미지에만 존재하는 정보를 포함한 문서를 만들어서,
텍스트만 추출 vs 이미지 포함 추출의 차이를 명확히 비교할 수 있게 한다.

사용법:
  pip install python-docx python-pptx matplotlib
  python create_test_files.py
"""

import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt


OUTPUT_DIR = Path(__file__).parent / "test_files"


def get_cjk_font():
    """시스템에서 한글 지원 폰트를 찾는다."""
    candidates = [
        "AppleGothic",
        "Apple SD Gothic Neo",
        "Malgun Gothic",
        "NanumGothic",
        "Noto Sans CJK KR",
        "Noto Sans KR",
    ]
    available = {f.name for f in fm.fontManager.ttflist}
    for name in candidates:
        if name in available:
            return name
    return None


def setup_matplotlib():
    cjk_font = get_cjk_font()
    if cjk_font:
        plt.rcParams["font.family"] = cjk_font
    plt.rcParams["axes.unicode_minus"] = False


def create_chart_quarterly_revenue(path: str):
    """분기별 매출 차트 — 이 수치는 이미지에만 존재"""
    setup_matplotlib()
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    revenue = [120, 145, 198, 176]  # 억원 — 이 숫자가 핵심!
    colors = ["#4285F4", "#34A853", "#FBBC04", "#EA4335"]

    fig, ax = plt.subplots(figsize=(6, 4))
    bars = ax.bar(quarters, revenue, color=colors, width=0.6)
    for bar, val in zip(bars, revenue):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 3,
                f"{val}", ha="center", va="bottom", fontweight="bold", fontsize=12)

    ax.set_title("2024 Quarterly Revenue", fontsize=14, fontweight="bold")
    ax.set_ylabel("Revenue (100M KRW)", fontsize=11)
    ax.set_ylim(0, 230)
    ax.spines[["top", "right"]].set_visible(False)
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()


def create_chart_market_share(path: str):
    """시장 점유율 파이 차트 — 이 비율은 이미지에만 존재"""
    setup_matplotlib()
    labels = ["Company A (42%)", "Company B (28%)", "Company C (18%)", "Others (12%)"]
    sizes = [42, 28, 18, 12]  # 이 수치가 핵심!
    colors = ["#4285F4", "#EA4335", "#FBBC04", "#34A853"]

    fig, ax = plt.subplots(figsize=(5, 5))
    ax.pie(sizes, labels=labels, colors=colors, autopct="%1.0f%%",
           startangle=90, textprops={"fontsize": 11})
    ax.set_title("2024 Market Share", fontsize=14, fontweight="bold")
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()


def create_chart_growth_trend(path: str):
    """연도별 성장 추이 — 이 트렌드는 이미지에만 존재"""
    setup_matplotlib()
    years = [2019, 2020, 2021, 2022, 2023, 2024]
    revenue = [85, 72, 110, 135, 162, 198]   # 이 수치가 핵심!
    profit = [12, -5, 18, 25, 32, 41]         # 이 수치도 핵심!

    fig, ax1 = plt.subplots(figsize=(7, 4))
    ax1.plot(years, revenue, "o-", color="#4285F4", linewidth=2,
             markersize=8, label="Revenue")
    ax1.set_ylabel("Revenue (100M KRW)", color="#4285F4", fontsize=11)

    ax2 = ax1.twinx()
    ax2.bar(years, profit, alpha=0.3, color="#34A853", width=0.5, label="Net Profit")
    ax2.set_ylabel("Net Profit (100M KRW)", color="#34A853", fontsize=11)

    ax1.set_title("Revenue & Profit Trend (2019-2024)",
                  fontsize=14, fontweight="bold")
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc="upper left")

    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()


def create_test_docx():
    """
    테스트용 DOCX 생성.
    텍스트에는 일반적인 내용만 있고,
    핵심 수치(매출, 점유율)는 차트 이미지에만 존재한다.
    """
    doc = Document()

    doc.add_heading("Annual Business Report 2024", level=1)
    doc.add_paragraph(
        "This report summarizes the business performance for fiscal year 2024. "
        "The company has shown strong growth across all segments. "
        "Detailed financial figures are presented in the charts below."
    )

    # 차트 1: 분기별 매출 (수치는 이미지에만 존재!)
    doc.add_heading("Quarterly Revenue Overview", level=2)
    doc.add_paragraph(
        "The following chart shows our quarterly revenue performance. "
        "We achieved record-breaking results in Q3."
        # 주의: 텍스트에 구체적 수치 없음 → 이미지에서만 파악 가능
    )
    chart_path = str(OUTPUT_DIR / "chart_revenue.png")
    create_chart_quarterly_revenue(chart_path)
    doc.add_picture(chart_path, width=Inches(5))

    # 차트 2: 시장 점유율 (비율은 이미지에만 존재!)
    doc.add_heading("Market Position", level=2)
    doc.add_paragraph(
        "Our market share has expanded significantly this year. "
        "The pie chart below illustrates our competitive position."
        # 주의: 텍스트에 구체적 비율 없음
    )
    chart_path2 = str(OUTPUT_DIR / "chart_market.png")
    create_chart_market_share(chart_path2)
    doc.add_picture(chart_path2, width=Inches(4))

    # 차트 3: 성장 추이 (트렌드는 이미지에만 존재!)
    doc.add_heading("Historical Growth", level=2)
    doc.add_paragraph(
        "The company has been on a consistent growth trajectory since 2019, "
        "with the exception of the pandemic year. "
        "Revenue and profit trends are shown below."
        # 주의: 텍스트에 연도별 수치 없음
    )
    chart_path3 = str(OUTPUT_DIR / "chart_growth.png")
    create_chart_growth_trend(chart_path3)
    doc.add_picture(chart_path3, width=Inches(5.5))

    doc.add_heading("Conclusion", level=2)
    doc.add_paragraph(
        "Looking ahead to 2025, we remain optimistic about continued growth "
        "driven by market expansion and operational efficiency improvements."
    )

    out_path = OUTPUT_DIR / "test_report.docx"
    doc.save(str(out_path))
    print(f"Created: {out_path}")
    return str(out_path)


def create_test_pptx():
    """
    테스트용 PPTX 생성.
    슬라이드 텍스트에는 일반 설명만 있고,
    핵심 데이터는 차트 이미지에만 존재한다.
    """
    prs = Presentation()

    # 슬라이드 1: 타이틀
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "2024 Business Review"
    slide.placeholders[1].text = "Annual Performance Summary"

    # 슬라이드 2: 매출 차트
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Quarterly Revenue - Record-breaking Q3 performance"
    p.font.size = Pt(24)
    p.font.bold = True

    chart_path = str(OUTPUT_DIR / "chart_revenue.png")
    slide.shapes.add_picture(
        chart_path, PptxInches(1.5), PptxInches(1.5),
        PptxInches(7), PptxInches(4.5))

    # 슬라이드 3: 점유율 차트
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Share - Expanded competitive position"
    p.font.size = Pt(24)
    p.font.bold = True

    chart_path2 = str(OUTPUT_DIR / "chart_market.png")
    slide.shapes.add_picture(
        chart_path2, PptxInches(2.5), PptxInches(1.2),
        PptxInches(5), PptxInches(5))

    # 슬라이드 4: 성장 추이 차트
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(
        PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Growth Trend - Consistent upward trajectory since 2019"
    p.font.size = Pt(24)
    p.font.bold = True

    chart_path3 = str(OUTPUT_DIR / "chart_growth.png")
    slide.shapes.add_picture(
        chart_path3, PptxInches(1), PptxInches(1.5),
        PptxInches(8), PptxInches(5))

    out_path = OUTPUT_DIR / "test_presentation.pptx"
    prs.save(str(out_path))
    print(f"Created: {out_path}")
    return str(out_path)


if __name__ == "__main__":
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    print("Generating test files...\n")
    create_test_docx()
    create_test_pptx()
    print(f"\nDone! Files are in: {OUTPUT_DIR}/")
    print("\nTest with:")
    print(f"  python test_image_understanding.py {OUTPUT_DIR}/test_report.docx")
    print(f"  python test_image_understanding.py {OUTPUT_DIR}/test_presentation.pptx")
